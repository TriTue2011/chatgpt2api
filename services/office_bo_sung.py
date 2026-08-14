"""Tám việc tài liệu mà `officecli` chưa làm được — soi từ repo Quiz99/officeCliMCP.

Vì sao KHÔNG cắm repo đó vào làm một máy chủ MCP riêng
------------------------------------------------------
`officeCliMCP` là một máy chủ MCP độc lập, viết sẵn để chạy cạnh dự án này (mã của
nó có đoạn ánh xạ `/app/data/images/docs/` sang `/app/data/c2a-docs`). Nhưng 19
trong 23 tool của nó trùng việc với thứ đã có ở đây:

    read_document/read_sheet/read_slides  ≈ officecli view/get + extract_markdown
    extract_images                        ≈ pdf_images.extract_office_images
    insert_content/delete_content         ≈ officecli add / remove
    edit_sheet                            ≈ officecli set
    create_document                       ≈ officecli create
    convert_document                      ≈ pdf_to_word / pdf_to_excel / markitdown
    summarize_document                    ≈ pdf_intent.summarize_pdf
    list_files                            ≈ officecli list_files

Dựng thêm một tiến trình, một cổng, một bộ phụ thuộc là đắt; và mỗi tool MCP còn
tốn context MỖI LƯỢT chat (xem ghi chú `skills.py` trong CLAUDE.md). Nên bù đúng
những việc thiếu, viết thẳng vào đây:

    1. so_sanh            so hai tài liệu                    (diff_documents)
    2. thong_ke_bang      thống kê bảng tính                 (analyze_data)
    3. tim_thay_the       tìm & thay thế toàn tài liệu       (replace_all)
    4. tao_slide          dàn ý → .pptx                      (create_slides)
    5. cat_theo_tieu_de   cắt tài liệu theo tiêu đề          (split_document)
    6. noi_tep            nối nhiều tệp thành một            (merge_documents)
    7. doc_thong_tin      tác giả / số trang / số sheet…     (extract_metadata)
    8. tao_bao_cao        Excel + mẫu Markdown → .docx       (create_report)

Chỗ bản này LÀM KHÁC repo, có lý do
------------------------------------
1. **Khoá trong thư mục làm việc.** `_resolve_path` của repo chỉ ghép tiền tố rồi
   ánh xạ vài đường dẫn, KHÔNG kiểm gì — `diff_documents("/etc/passwd", …)` là đọc
   được tệp đó. Ở đây dùng `officecli.resolve_path`, nó `resolve()` rồi đòi đường
   dẫn phải nằm trong workspace, nên `..` và symlink đều không lách được. Bot đọc
   tài liệu người lạ gửi tới, mà tài liệu có thể chứa câu ra lệnh.
2. **Đọc tài liệu bằng đường sẵn có.** Repo đọc .docx bằng cách nối `p.text` của
   từng paragraph — MẤT HẾT bảng. Ở đây gọi `pdf_intent.extract_markdown`
   (pdf-inspector → PyMuPDF → markitdown) nên bảng còn nguyên, và PDF scan còn
   được OCR. So sánh hai bản hợp đồng mà rơi mất bảng thì so sánh để làm gì.
3. **Không raise ra ngoài.** Mọi hàm trả `{ok, ...}`; đây là việc phụ của luồng
   chat, hỏng thì báo lại chứ không được làm gãy lượt trả lời.
"""
from __future__ import annotations

import difflib
import logging
import re
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

#: Trần độ dài báo cáo trả về bot. Vượt trần thì cắt và NÓI LÀ ĐÃ CẮT — cắt âm
#: thầm là thứ đã làm mất cả buổi để lần ra ở `ask_choices`.
MAX_KY_TU = 12000


def _duong_dan(p: str) -> Path:
    from services import officecli
    return officecli.resolve_path(p, must_exist=True)


def _doc_chu(p: Path) -> str:
    """Tài liệu → chữ, giữ bảng. Dùng đúng đường trích xuất của dự án."""
    duoi = p.suffix.lower()
    if duoi in {".md", ".txt", ".csv", ".json"}:
        return p.read_text("utf-8", errors="ignore")
    from services import pdf_intent
    return pdf_intent.extract_markdown(str(p)) or ""


def _cat(s: str) -> tuple[str, bool]:
    return (s[:MAX_KY_TU], True) if len(s) > MAX_KY_TU else (s, False)


# ── 1. So sánh hai tài liệu ─────────────────────────────────────────────────

def so_sanh(tep_a: str, tep_b: str, *, dinh_dang: str = "markdown") -> dict[str, Any]:
    """So hai tài liệu, trả báo cáo khác biệt. dinh_dang: markdown | unified."""
    try:
        a, b = _duong_dan(tep_a), _duong_dan(tep_b)
    except (ValueError, FileNotFoundError, OSError) as exc:
        return {"ok": False, "error": str(exc)}
    try:
        ta, tb = _doc_chu(a), _doc_chu(b)
    except Exception as exc:
        return {"ok": False, "error": f"không đọc được nội dung: {str(exc)[:150]}"}

    da, db = ta.splitlines(), tb.splitlines()
    if da == db:
        return {"ok": True, "giong_nhau": True,
                "bao_cao": f"✅ {a.name} và {b.name} giống nhau.",
                "them": 0, "bo": 0, "sua": 0}

    if dinh_dang == "unified":
        bc = "\n".join(difflib.unified_diff(da, db, fromfile=a.name, tofile=b.name,
                                            lineterm=""))
        them = sum(1 for d in bc.splitlines() if d.startswith("+") and not d.startswith("+++"))
        bo = sum(1 for d in bc.splitlines() if d.startswith("-") and not d.startswith("---"))
        bc, bi_cat = _cat(bc)
        return {"ok": True, "giong_nhau": False, "bao_cao": bc,
                "them": them, "bo": bo, "sua": 0, "bi_cat": bi_cat}

    them = bo = sua = 0
    khoi: list[str] = []
    for tag, i1, i2, j1, j2 in difflib.SequenceMatcher(None, da, db).get_opcodes():
        if tag == "equal":
            continue
        if tag == "replace":
            sua += 1
            khoi.append(f"\n**🔄 Sửa (dòng {i1 + 1})**\n")
            khoi += [f"- {x}" for x in da[i1:i2]]
            khoi += [f"+ {x}" for x in db[j1:j2]]
        elif tag == "delete":
            bo += i2 - i1
            khoi.append(f"\n**❌ Bỏ (dòng {i1 + 1})**\n")
            khoi += [f"- {x}" for x in da[i1:i2]]
        elif tag == "insert":
            them += j2 - j1
            khoi.append(f"\n**✅ Thêm (dòng {j1 + 1})**\n")
            khoi += [f"+ {x}" for x in db[j1:j2]]
    dau = (f"# So sánh {a.name} ↔ {b.name}\n\n"
           f"- Thêm: {them} dòng\n- Bỏ: {bo} dòng\n- Sửa: {sua} khối\n")
    bc, bi_cat = _cat(dau + "\n".join(khoi))
    return {"ok": True, "giong_nhau": False, "bao_cao": bc,
            "them": them, "bo": bo, "sua": sua, "bi_cat": bi_cat}


# ── 2. Thống kê bảng Excel ──────────────────────────────────────────────────

def thong_ke_bang(tep: str, *, sheet: str = "") -> dict[str, Any]:
    """Thống kê nhanh một bảng tính: số dòng/cột, thiếu dữ liệu, min/max/trung bình."""
    try:
        p = _duong_dan(tep)
    except (ValueError, FileNotFoundError, OSError) as exc:
        return {"ok": False, "error": str(exc)}
    if p.suffix.lower() not in {".xlsx", ".xls", ".xlsm", ".csv"}:
        return {"ok": False, "error": f"chỉ thống kê được bảng tính, không phải {p.suffix}"}
    try:
        import pandas as pd
        if p.suffix.lower() == ".csv":
            bang = {"CSV": pd.read_csv(str(p))}
        else:
            bang = pd.read_excel(str(p), sheet_name=(sheet or None))
            if not isinstance(bang, dict):
                bang = {sheet or "Sheet1": bang}
    except Exception as exc:
        return {"ok": False, "error": f"không đọc được bảng: {str(exc)[:150]}"}

    dong: list[str] = [f"# Thống kê {p.name}"]
    tong_dong = 0
    for ten, df in bang.items():
        if df is None or df.empty:
            dong.append(f"\n## {ten}\n(trống)")
            continue
        tong_dong += len(df)
        dong.append(f"\n## {ten}\n\n{len(df)} dòng × {len(df.columns)} cột")
        thieu = {c: int(df[c].isna().sum()) for c in df.columns if df[c].isna().any()}
        if thieu:
            dong.append("\n**Thiếu dữ liệu:** "
                        + ", ".join(f"{c} ({n})" for c, n in thieu.items()))
        so = df.select_dtypes("number")
        if not so.empty:
            dong.append("\n| Cột | Nhỏ nhất | Lớn nhất | Trung bình | Tổng |")
            dong.append("|---|---|---|---|---|")
            for c in so.columns:
                v = so[c].dropna()
                if v.empty:
                    continue
                dong.append(f"| {c} | {v.min():g} | {v.max():g} | "
                            f"{v.mean():.2f} | {v.sum():g} |")
        chu = [c for c in df.columns if c not in so.columns]
        for c in chu[:3]:      # ba cột chữ đầu, đủ để nhận ra dữ liệu là gì
            top = df[c].astype(str).value_counts().head(3)
            if len(top):
                dong.append(f"\n**{c}** hay gặp: "
                            + ", ".join(f"{k} ({v})" for k, v in top.items()))
    bc, bi_cat = _cat("\n".join(dong))
    return {"ok": True, "bao_cao": bc, "so_sheet": len(bang),
            "tong_dong": tong_dong, "bi_cat": bi_cat}


# ── 3. Tìm & thay thế toàn tài liệu ────────────────────────────────────────
# `officecli set` chỉ sửa được ĐÚNG một phần tử đã biết đường dẫn. Muốn đổi một
# cụm từ nằm rải khắp tài liệu thì phải tra từng chỗ — việc thường gặp nhất
# (đổi tên, đổi năm học) lại là việc khó nhất.

def tim_thay_the(tep: str, tim: str, thay: str, *,
                 tat_ca: bool = True) -> dict[str, Any]:
    """Đổi `tim` thành `thay` trong .docx / .xlsx. Ghi thẳng vào tệp đó."""
    if not str(tim or ""):
        return {"ok": False, "error": "chưa cho biết cần tìm cụm gì"}
    try:
        p = _duong_dan(tep)
    except (ValueError, FileNotFoundError, OSError) as exc:
        return {"ok": False, "error": str(exc)}
    duoi = p.suffix.lower()
    try:
        if duoi == ".docx":
            n = _thay_docx(p, tim, thay, tat_ca)
        elif duoi in {".xlsx", ".xlsm"}:
            n = _thay_xlsx(p, tim, thay, tat_ca)
        else:
            return {"ok": False, "error": f"chỉ đổi được .docx và .xlsx, không phải {duoi}"}
    except Exception as exc:
        logger.warning("tim_thay_the %s lỗi: %s", p.name, str(exc)[:150])
        return {"ok": False, "error": f"không sửa được: {str(exc)[:150]}"}
    if not n:
        return {"ok": True, "so_lan": 0,
                "text": f"Không thấy “{tim}” trong {p.name} ạ."}
    return {"ok": True, "so_lan": n,
            "text": f"Đã đổi {n} chỗ trong {p.name} ạ."}


def _thay_trong_doan(doan: Any, tim: str, thay: str, con_lai: list[int]) -> int:
    """Đổi trong một paragraph, GIỮ ĐỊNH DẠNG của run đầu tiên khớp.

    Vì sao không gán thẳng `doan.text = ...`: làm thế là xoá sạch mọi run, tức
    mất hết đậm/nghiêng/màu của cả đoạn. Ở đây chỉ sửa text của từng run.
    """
    if tim not in doan.text:
        return 0
    dem = 0
    for run in doan.runs:
        if tim not in run.text or con_lai[0] == 0:
            continue
        if con_lai[0] < 0:                       # -1 = đổi tất cả
            so = run.text.count(tim)
            run.text = run.text.replace(tim, thay)
        else:
            so = min(run.text.count(tim), con_lai[0])
            run.text = run.text.replace(tim, thay, so)
            con_lai[0] -= so
        dem += so
    # Cụm bị cắt ngang giữa hai run thì vòng trên không thấy. Gộp cả đoạn vào
    # run đầu — mất định dạng riêng của các run sau, nhưng đúng nội dung; thà
    # vậy còn hơn im lặng không đổi gì.
    if dem == 0 and doan.runs:
        doan.runs[0].text = doan.text.replace(tim, thay)
        for r in doan.runs[1:]:
            r.text = ""
        dem = doan.text.count(tim) or 1
    return dem


def _thay_docx(p: Path, tim: str, thay: str, tat_ca: bool) -> int:
    import docx
    d = docx.Document(str(p))
    con_lai = [-1 if tat_ca else 1]
    dem = 0
    for doan in d.paragraphs:
        dem += _thay_trong_doan(doan, tim, thay, con_lai)
        if con_lai[0] == 0:
            break
    if con_lai[0] != 0:
        for bang in d.tables:          # bảng: repo bỏ qua hoàn toàn
            for hang in bang.rows:
                for o in hang.cells:
                    for doan in o.paragraphs:
                        dem += _thay_trong_doan(doan, tim, thay, con_lai)
    if dem:
        d.save(str(p))
    return dem


def _thay_xlsx(p: Path, tim: str, thay: str, tat_ca: bool) -> int:
    import openpyxl
    wb = openpyxl.load_workbook(str(p))
    dem = 0
    for ws in wb.worksheets:
        for hang in ws.iter_rows():
            for o in hang:
                if isinstance(o.value, str) and tim in o.value:
                    so = o.value.count(tim)
                    o.value = o.value.replace(tim, thay) if tat_ca \
                        else o.value.replace(tim, thay, 1)
                    dem += so if tat_ca else 1
                    if not tat_ca:
                        wb.save(str(p))
                        return dem
    if dem:
        wb.save(str(p))
    return dem


# ── 4. Tạo slide PowerPoint từ dàn ý ───────────────────────────────────────
# Dự án chưa tạo được .pptx ở bất kỳ đâu (soi cả `services/`: không nơi nào
# import pptx để GHI). Với đường dạy học thì đây là việc đáng có nhất trong repo.
#
# Bản đầu chỉ đổ chữ vào bố cục mặc định của python-pptx: mọi slide một kiểu,
# màu và cỡ chữ theo mẫu trắng trơn, bảng số liệu thành mấy dòng gạch đầu dòng.
# Bản này chọn BỐ CỤC THEO LOẠI NỘI DUNG và lấy màu/cỡ chữ từ MỘT bộ chủ đề:
#
#     `#` không có ý nào, ở đầu   → slide bìa (tựa lớn + phụ đề)
#     `#` không có ý nào, ở giữa  → slide phân mục
#     có dòng `| a | b |`         → BẢNG GỐC của PowerPoint (sửa được trong app)
#     có `[biểu đồ]` trước bảng   → BIỂU ĐỒ GỐC (chuột phải → Edit Data ra số)
#     còn lại                     → gạch đầu dòng, thụt lề thành cấp thật
#
# Vẫn thuần python-pptx, không thêm thư viện nào. Cố ý KHÔNG đi đường
# SVG → DrawingML như ppt-master: đường đó đẹp hơn nhưng phải tự bảo trì một
# tầng chuyển đổi, mà đây chỉ là việc phụ của bot.

_RE_TIEU_DE = re.compile(r"^(#{1,3})\s+(.*)$")
#: Nhóm 1 giữ phần thụt lề — hai dấu cách là một cấp (xem `_cap_thut_le`).
_RE_GACH_DAU = re.compile(r"^(\s*)[-*•]\s+(.*)$")
#: Một hàng bảng Markdown: `| a | b |`.
_RE_HANG_BANG = re.compile(r"^\s*\|(.+)\|\s*$")
#: Dòng ngăn dưới hàng tiêu đề bảng: `| --- | :--: |` — bỏ, không phải dữ liệu.
_RE_NGAN_BANG = re.compile(r"^\s*\|[\s:|+-]+\|\s*$")
#: Mốc xin biểu đồ. Cố ý bắt người soạn NÓI RÕ thay vì tự đoán theo dữ liệu:
#: bảng số liệu tự biến thành biểu đồ là kiểu "thông minh" gây bất ngờ, mà mất
#: luôn phần đọc được từng con số.
_RE_MOC_BIEU_DO = re.compile(
    r"^\s*[\[(]\s*(?:bi[eể]u\s*đ[oồ]|chart|graph)\s*[\])]\s*$", re.I)

#: Bộ chủ đề. Cố ý tránh gradient tím/xanh-tím — đó là dấu hiệu dễ nhận nhất
#: của giao diện do AI sinh, và slide bài giảng thì cần trông có chủ ý.
#: `font` để Calibri: nó có sẵn trên mọi máy Office và đủ dấu tiếng Việt. Font
#: lạ mà máy người nhận không có thì PowerPoint thay bừa, slide vỡ hết chữ.
_CHU_DE: dict[str, dict[str, Any]] = {
    # Trắng sạch — gần nhất với đầu ra cũ, nên để làm mặc định: bản nâng cấp
    # không được làm slide người ta đã quen bỗng đổi màu.
    "trang-sach": {"nen": "FFFFFF", "tua": "1A2433", "chu": "333F4F",
                   "nhan": "B4553C", "chu_tren_nhan": "FFFFFF"},
    "xanh-dam":   {"nen": "0F243E", "tua": "F2F5F7", "chu": "C9D4DE",
                   "nhan": "E0A458", "chu_tren_nhan": "1A2433"},
    "am":         {"nen": "1C1B1A", "tua": "F5F0E8", "chu": "CFC7BC",
                   "nhan": "7A9E7E", "chu_tren_nhan": "1C1B1A"},
}
_CHU_DE_MAC_DINH = "trang-sach"
_FONT = "Calibri"
#: Cỡ chữ (pt). Ý cấp sâu nhỏ dần nhưng có sàn — chữ dưới 14pt thì cuối phòng
#: học không đọc được, mà slide bài giảng là để cả phòng đọc.
_CO_TUA_BIA, _CO_PHU_BIA, _CO_TUA, _CO_Y, _CO_Y_SAN, _CO_BANG = 40, 20, 32, 20, 14, 14


def _cap_thut_le(khoang_trang: str) -> int:
    """Phần thụt lề → cấp gạch đầu dòng (0–3). Tab tính bằng bốn dấu cách."""
    n = len(khoang_trang.replace("\t", "    "))
    return max(0, min(3, n // 2))


#: Ba KHUÔN số được chấp. Cố ý khớp khuôn thay vì bóc lấy chữ số: hàm này quyết
#: định một cột bảng có phải số liệu hay không, nên "bóc được vài chữ số rồi trả
#: về số" chính là cách vẽ ra một biểu đồ vô nghĩa từ dữ liệu rác.
#: Kiểu Việt — chấm phân nhóm nghìn, phẩy thập phân: 1.200 · 1.200,5 · 1.200.000
_RE_SO_VIET = re.compile(r"[+-]?\d{1,3}(?:\.\d{3})+(?:,\d+)?")
#: Kiểu Anh — ngược lại: 1,200 · 1,200.5
_RE_SO_ANH = re.compile(r"[+-]?\d{1,3}(?:,\d{3})+(?:\.\d+)?")
#: Không nhóm nghìn, nhiều nhất một dấu thập phân: 8 · 1200 · 1,5 · 1.5 · -2,5
_RE_SO_TRON = re.compile(r"[+-]?\d+(?:[.,]\d+)?")


def _so(chu: str) -> float | None:
    """Chuỗi → số, hiểu cả lối viết Việt lẫn Anh. None nếu không phải số.

    Phân biệt hai lối bằng số chữ số sau dấu: đúng BA chữ số thì dấu đó phân
    nhóm nghìn, khác ba thì nó là dấu thập phân.

        1.200 → 1200      1.200,5 → 1200,5      (Việt)
        1,200 → 1200      1,200.5 → 1200,5      (Anh)
        1,5 · 1.5 → 1,5   1 200 → 1200    12% → 12

    Cùng một dấu vừa phân nhóm vừa làm thập phân là SAI khuôn → trả None chứ
    không bóc chữ số ra đoán bừa: `1.2.3,4,5` phải bị loại.

    Chỗ CỐ Ý bỏ: ai viết `1.200` mà thật sự muốn một phẩy hai thì bị đọc thành
    một nghìn hai. Không có cách nào phân biệt hai ý đó từ chuỗi, và đây là bot
    tiếng Việt nên chọn nghĩa người Việt hay dùng hơn.
    """
    s = str(chu or "").strip().replace(" ", "").replace(" ", "")
    s = s.replace("%", "")
    if not s:
        return None
    if _RE_SO_VIET.fullmatch(s):
        s = s.replace(".", "").replace(",", ".")
    elif _RE_SO_ANH.fullmatch(s):
        s = s.replace(",", "")
    elif _RE_SO_TRON.fullmatch(s):
        s = s.replace(",", ".")
    else:
        return None
    try:
        return float(s)
    except ValueError:                       # khuôn đã chặn; đây là lưới cuối
        return None


#: Dạng MƠ HỒ: một dấu phân cách, đúng ba chữ số sau nó, không dấu nào khác.
#: `1.200` có thể là một nghìn hai (nhóm nghìn) hoặc một phẩy hai (thập phân).
_RE_SO_MO_HO = re.compile(r"[+-]?\d{1,3}[.,]\d{3}")


def _cot_so(cot: list[str]) -> list[float] | None:
    """Cột bảng → dãy số, dùng CẢ CỘT để gỡ chỗ mơ hồ. None nếu có ô không phải số.

    Vì sao không đọc từng ô bằng `_so`: một mình thì `6.750` mơ hồ thật, nhưng cả
    cột thì không — các ô trong một cột là cùng một loại đại lượng, nên phải đọc
    sao cho chúng cùng cỡ. Điểm trung bình `6.750 / 7,25 / 8` đọc kiểu nhóm nghìn
    ra 6750 / 7,25 / 8: một cột cao gấp gần nghìn lần hai cột kia, tức là đọc sai.
    Đọc kiểu thập phân ra 6,75 / 7,25 / 8 mới đúng ý người viết. Đây là lỗi đo
    được trên bộ slide thử, không phải lo xa.

    Cách gỡ: lấy các ô KHÔNG mơ hồ làm mốc, rồi mỗi ô mơ hồ chọn cách đọc cho giá
    trị gần cỡ mốc hơn. Cả cột đều mơ hồ (`1.200 / 1.450`) thì giữ nghĩa nhóm
    nghìn — đó là lối người Việt hay viết hơn.
    """
    tho = [str(o or "").strip().replace(" ", "").replace("%", "") for o in cot]
    gia_tri = [_so(o) for o in tho]
    if not gia_tri or any(v is None for v in gia_tri):
        return None
    mo_ho = [i for i, o in enumerate(tho) if _RE_SO_MO_HO.fullmatch(o)]
    moc = [abs(v) for i, v in enumerate(gia_tri)
           if i not in mo_ho and v and abs(v) > 0]
    if not mo_ho or not moc:
        return [float(v) for v in gia_tri]      # type: ignore[arg-type]
    giua = sorted(moc)[len(moc) // 2]
    for i in mo_ho:
        nhom = float(gia_tri[i])                # cách đọc mặc định: nhóm nghìn
        thap_phan = nhom / 1000.0               # cùng chữ số, dời dấu phẩy
        if abs(thap_phan - giua) < abs(nhom - giua):
            gia_tri[i] = thap_phan
    return [float(v) for v in gia_tri]          # type: ignore[arg-type]


def _tach_dan_y(noi_dung: str) -> list[dict[str, Any]]:
    """Dàn ý Markdown → danh sách mô tả slide.

    Mỗi mục: ``{"tua", "cap", "y": [(cấp, chữ)], "bang": [[ô]], "bieu_do": bool}``.
    Chữ đứng trước tiêu đề đầu tiên vẫn thành một slide mở đầu — giữ đúng hành vi
    cũ, vì có người gửi dàn ý không đánh tiêu đề nào.
    """
    cac_slide: list[dict[str, Any]] = []

    def slide_hien_tai() -> dict[str, Any]:
        if not cac_slide:
            cac_slide.append({"tua": "", "cap": 1, "y": [], "bang": [],
                              "bieu_do": False})
        return cac_slide[-1]

    xin_bieu_do = False
    for dong in noi_dung.splitlines():
        m = _RE_TIEU_DE.match(dong.strip())
        if m:
            cac_slide.append({"tua": m.group(2).strip(), "cap": len(m.group(1)),
                              "y": [], "bang": [], "bieu_do": False})
            xin_bieu_do = False
            continue
        if _RE_MOC_BIEU_DO.match(dong):
            xin_bieu_do = True
            continue
        if _RE_NGAN_BANG.match(dong):
            continue
        hb = _RE_HANG_BANG.match(dong)
        if hb:
            s = slide_hien_tai()
            s["bang"].append([o.strip() for o in hb.group(1).split("|")])
            if xin_bieu_do:
                s["bieu_do"] = True
            continue
        g = _RE_GACH_DAU.match(dong)
        if g:
            slide_hien_tai()["y"].append((_cap_thut_le(g.group(1)),
                                          g.group(2).strip()))
            continue
        chu = dong.strip()
        if chu:
            slide_hien_tai()["y"].append((0, chu))
    return cac_slide


def _co_the_ve_bieu_do(bang: list[list[str]]) -> bool:
    """Bảng có vẽ được biểu đồ cột không: ≥2 cột, ≥2 hàng dữ liệu, và cột cuối
    toàn số. Không đạt thì vẽ bảng — thà đúng còn hơn ra một biểu đồ rỗng."""
    if len(bang) < 3 or len(bang[0]) < 2:
        return False
    return _cot_so([h[-1] for h in bang[1:] if len(h) >= 2]) is not None


def tao_slide(dan_y: str, ten_tep: str, chu_de: str = "") -> dict[str, Any]:
    """Dàn ý Markdown → .pptx. Mỗi tiêu đề `#`/`##` là một slide, gạch đầu dòng
    là các ý trong slide đó.

    `chu_de`: `trang-sach` (mặc định) | `xanh-dam` | `am`. Tên lạ thì dùng mặc
    định chứ không báo hỏng — tạo slide là việc phụ của lượt chat, không được
    gãy chỉ vì gõ sai tên màu.

    Bố cục chọn theo NỘI DUNG, không phải một khuôn cho mọi slide:
      * `#` không có ý nào, đứng đầu dàn ý → slide bìa;
      * `#` không có ý nào, ở giữa → slide phân mục;
      * có dòng `| a | b |` → bảng gốc của PowerPoint;
      * có dòng `[biểu đồ]` trước bảng → biểu đồ cột gốc (cần ≥2 hàng dữ liệu và
        cột cuối toàn số, không đạt thì vẽ bảng);
      * còn lại → gạch đầu dòng, thụt lề hai dấu cách thành một cấp.

    Trả thêm `so_bang`, `so_bieu_do`, `chu_de` so với bản cũ; các khoá cũ
    (`ok`, `duong_dan`, `ten`, `so_slide`) giữ nguyên.
    """
    noi_dung = str(dan_y or "").strip()
    if not noi_dung:
        return {"ok": False, "error": "chưa có dàn ý"}
    ten = re.sub(r"[^\w.() \-]+", "_", Path(str(ten_tep or "")).name).strip()
    if not ten:
        ten = "slide.pptx"
    if not ten.lower().endswith(".pptx"):
        ten += ".pptx"
    try:
        from services import officecli
        p = officecli.resolve_path(ten)
    except (ValueError, OSError) as exc:
        return {"ok": False, "error": str(exc)}

    cac_slide = _tach_dan_y(noi_dung)
    if not cac_slide:
        return {"ok": False, "error": "dàn ý không có tiêu đề hay ý nào"}
    ten_cd = str(chu_de or "").strip().lower() or _CHU_DE_MAC_DINH
    if ten_cd not in _CHU_DE:               # tên chủ đề lạ → dùng mặc định,
        ten_cd = _CHU_DE_MAC_DINH           # không báo hỏng vì đây là việc phụ
    cd = _CHU_DE[ten_cd]

    so_bang = so_bieu_do = 0
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt

        def mau(khoa: str):
            return RGBColor.from_string(cd[khoa])

        def dat_chu(khung, co, khoa_mau: str, *, dam=False) -> None:
            """Ép font/cỡ/màu lên TỪNG run. Đặt ở cấp paragraph là không đủ:
            run sinh sau (add_paragraph) không thừa hưởng, nên slide ra nửa
            theo chủ đề nửa theo mẫu."""
            for doan in khung.paragraphs:
                for r in doan.runs or [doan.add_run()]:
                    r.font.name, r.font.size = _FONT, Pt(co)
                    r.font.color.rgb, r.font.bold = mau(khoa_mau), dam

        pres = Presentation()
        # 16:9. Mẫu trần của python-pptx là 4:3 — khổ đó chiếu lên màn hình hay
        # máy chiếu bây giờ là hai vệt đen hai bên.
        pres.slide_width, pres.slide_height = Inches(13.333), Inches(7.5)
        rong = pres.slide_width

        for thu_tu, sl in enumerate(cac_slide):
            tua, cac_y, bang = sl["tua"], sl["y"], sl["bang"]
            ve_bieu_do = bool(sl["bieu_do"]) and _co_the_ve_bieu_do(bang)
            if bang:
                bo_cuc = 5                                   # Title Only
            elif not cac_y and sl["cap"] == 1:
                bo_cuc = 0 if thu_tu == 0 else 2             # bìa / phân mục
            else:
                bo_cuc = 1                                   # Title and Content
            s = pres.slides.add_slide(pres.slide_layouts[bo_cuc])
            s.background.fill.solid()
            s.background.fill.fore_color.rgb = mau("nen")
            if s.shapes.title is not None:
                s.shapes.title.text = tua or " "
                dat_chu(s.shapes.title.text_frame,
                        _CO_TUA_BIA if bo_cuc == 0 else _CO_TUA, "tua", dam=True)

            if bo_cuc in (0, 2):
                # Bìa và phân mục: phụ đề chỉ là gạch chân màu nhấn. Đủ để slide
                # trông có chủ ý mà không bịa thêm chữ người soạn không viết.
                phu = next((ph for ph in s.placeholders
                            if ph.placeholder_format.idx == 1), None)
                if phu is not None:
                    phu.text_frame.text = ""
                    dat_chu(phu.text_frame, _CO_PHU_BIA, "chu")
                from pptx.enum.shapes import MSO_SHAPE
                gach = s.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.9), Inches(4.6), Inches(2.2), Inches(0.06))
                gach.fill.solid()
                gach.fill.fore_color.rgb = mau("nhan")
                gach.line.fill.background()
                continue

            if bang and ve_bieu_do:
                so_bieu_do += 1
                _ve_bieu_do(s, bang, cd, rong)
                continue
            if bang:
                so_bang += 1
                _ve_bang(s, bang, cd, rong)
                continue

            khung = next((ph.text_frame for ph in s.placeholders
                          if ph.placeholder_format.idx == 1), None)
            if khung is None:
                continue
            for i, (cap, chu) in enumerate(cac_y):
                doan = khung.paragraphs[0] if i == 0 else khung.add_paragraph()
                doan.text, doan.level = chu, cap
            for doan in khung.paragraphs:
                cap = doan.level or 0
                for r in doan.runs:
                    r.font.name = _FONT
                    r.font.size = Pt(max(_CO_Y_SAN, _CO_Y - 2 * cap))
                    r.font.color.rgb = mau("tua" if cap == 0 else "chu")

        p.parent.mkdir(parents=True, exist_ok=True)
        pres.save(str(p))
    except Exception as exc:
        logger.warning("tao_slide %s lỗi: %s", ten, str(exc)[:150])
        return {"ok": False, "error": f"không tạo được slide: {str(exc)[:150]}"}
    return {"ok": True, "duong_dan": str(p), "ten": p.name,
            "so_slide": len(cac_slide), "so_bang": so_bang,
            "so_bieu_do": so_bieu_do, "chu_de": ten_cd}


def _ve_bang(s, bang: list[list[str]], cd: dict[str, Any], rong) -> None:
    """Bảng Markdown → BẢNG GỐC của PowerPoint (không phải ảnh): người nhận sửa
    được từng ô. Hàng đầu tô màu nhấn để đọc ra ngay đâu là tiêu đề cột."""
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
    so_cot = max(len(h) for h in bang)
    le, tren = Inches(0.7), Inches(1.9)
    cao = Inches(min(4.8, 0.42 * len(bang)))
    tb = s.shapes.add_table(len(bang), so_cot, le, tren,
                            rong - Inches(1.4), cao).table
    for i, hang in enumerate(bang):
        for j in range(so_cot):
            o = tb.cell(i, j)
            o.text = hang[j] if j < len(hang) else ""
            o.fill.solid()
            o.fill.fore_color.rgb = RGBColor.from_string(
                cd["nhan"] if i == 0 else cd["nen"])
            for doan in o.text_frame.paragraphs:
                for r in doan.runs:
                    r.font.name, r.font.size = _FONT, Pt(_CO_BANG)
                    r.font.bold = i == 0
                    r.font.color.rgb = RGBColor.from_string(
                        cd["chu_tren_nhan"] if i == 0 else cd["chu"])


def _ve_bieu_do(s, bang: list[list[str]], cd: dict[str, Any], rong) -> None:
    """Bảng số liệu → BIỂU ĐỒ CỘT GỐC. Dữ liệu nằm trong file, chuột phải →
    Edit Data là ra đúng mấy con số này — không phải ảnh chụp biểu đồ."""
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches, Pt
    nhan = [h[0] for h in bang[1:]]
    gia_tri = _cot_so([h[-1] for h in bang[1:]]) or []
    dl = CategoryChartData()
    dl.categories = nhan
    dl.add_series(bang[0][-1] or "Số liệu", tuple(gia_tri))
    kh = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.0),
                            Inches(1.8), rong - Inches(2.0), Inches(4.9), dl)
    ch = kh.chart
    ch.has_legend = False             # một chuỗi thì chú giải chỉ là chữ thừa
    # Cột và chữ trục cũng phải theo chủ đề. Bỏ qua thì biểu đồ ra màu xanh mặc
    # định của Office, lệch hẳn với phần còn lại của bộ slide.
    to = ch.plots[0].series[0].format.fill
    to.solid()
    to.fore_color.rgb = RGBColor.from_string(cd["nhan"])
    for truc in (ch.category_axis, ch.value_axis):
        chu = truc.tick_labels.font
        chu.name, chu.size = _FONT, Pt(_CO_BANG)
        chu.color.rgb = RGBColor.from_string(cd["chu"])


# ── 5. Cắt tài liệu theo tiêu đề ────────────────────────────────────────────
# Repo cắt bằng `content.split("\n## ")` — CỨNG ở cấp 2. Tài liệu chỉ dùng `#`
# (rất thường gặp với văn bản tiếng Việt một cấp) thì không cắt được mảnh nào,
# trả về đúng một tệp bằng cả bản gốc mà vẫn báo ok.

_RE_TD = re.compile(r"^(#{1,6})\s+(.+?)\s*$")


def _cap_tieu_de_nong_nhat(cac_dong: list[str]) -> int:
    """Cấp tiêu đề NÔNG nhất có trong tài liệu (1 = '#'). 0 = không có tiêu đề."""
    cap = 0
    for d in cac_dong:
        m = _RE_TD.match(d)
        if m:
            n = len(m.group(1))
            cap = n if cap == 0 else min(cap, n)
    return cap


def cat_theo_tieu_de(tep: str, *, cap: int = 0,
                     thu_muc: str = "") -> dict[str, Any]:
    """Cắt tài liệu thành nhiều tệp .md theo tiêu đề.

    `cap=0` → tự lấy cấp tiêu đề NÔNG nhất tài liệu đang dùng, nên tài liệu chỉ
    có `#` vẫn cắt được. Bảng trong .docx còn nguyên vì đọc qua
    `pdf_intent.extract_markdown`.
    """
    try:
        p = _duong_dan(tep)
    except (ValueError, FileNotFoundError, OSError) as exc:
        return {"ok": False, "error": str(exc)}
    try:
        noi_dung = _doc_chu(p)
    except Exception as exc:
        return {"ok": False, "error": f"không đọc được: {str(exc)[:150]}"}
    dong = noi_dung.splitlines()
    muc_cap = int(cap) if cap else _cap_tieu_de_nong_nhat(dong)
    if not muc_cap:
        return {"ok": False,
                "error": "tài liệu không có tiêu đề Markdown nào để cắt theo"}

    try:
        from services import officecli
        ra_dir = (officecli.resolve_path(thu_muc) if str(thu_muc or "").strip()
                  else officecli.workspace() / f"{p.stem}-cat")
        ra_dir.mkdir(parents=True, exist_ok=True)
    except (ValueError, OSError) as exc:
        return {"ok": False, "error": str(exc)}

    phan: list[tuple[str, list[str]]] = []
    for d in dong:
        m = _RE_TD.match(d)
        if m and len(m.group(1)) == muc_cap:
            phan.append((m.group(2).strip(), [d]))
        elif phan:
            phan[-1][1].append(d)
        elif d.strip():
            phan.append(("mo-dau", [d]))     # chữ trước tiêu đề đầu tiên
    if not phan:
        return {"ok": False, "error": "không tách được phần nào"}

    cac_tep: list[str] = []
    for i, (ten, than) in enumerate(phan, 1):
        an_toan = re.sub(r"[^\w \-]+", "", ten)[:50].strip().replace(" ", "-") or f"phan-{i}"
        q = ra_dir / f"{i:02d}-{an_toan}.md"
        q.write_text("\n".join(than).strip() + "\n", "utf-8")
        cac_tep.append(q.name)
    return {"ok": True, "so_phan": len(cac_tep), "cap": muc_cap,
            "thu_muc": str(ra_dir), "cac_tep": cac_tep}


# ── 6. Nối nhiều tệp thành một ─────────────────────────────────────────────
# KHÁC `office_merge` sẵn có: cái đó là trộn dữ liệu vào MẪU (mail-merge), còn
# đây là ghép nội dung nhiều tệp lại.
#
# Repo nối .docx bằng `new_para.text = para.text` — mất sạch đậm/nghiêng, và bỏ
# luôn mọi bảng. Nó còn `continue` im lặng khi tệp thiếu, nên nối 5 tệp mà ra 2
# thì người dùng không biết ba tệp nào đã rơi.

def noi_tep(cac_tep: list[str], ten_ra: str, *,
            dinh_dang: str = "docx") -> dict[str, Any]:
    """Nối nhiều tài liệu thành một. dinh_dang: docx | md."""
    ds = [str(x) for x in (cac_tep or []) if str(x or "").strip()]
    if len(ds) < 2:
        return {"ok": False, "error": "cần ít nhất hai tệp để nối"}
    try:
        from services import officecli
        ra = officecli.resolve_path(str(ten_ra or "").strip() or "noi.docx")
    except (ValueError, OSError) as exc:
        return {"ok": False, "error": str(exc)}

    hop_le: list[Path] = []
    bo_qua: list[str] = []
    for x in ds:
        try:
            hop_le.append(_duong_dan(x))
        except (ValueError, FileNotFoundError, OSError) as exc:
            bo_qua.append(f"{x} ({str(exc)[:60]})")
    if len(hop_le) < 2:
        return {"ok": False, "error": "không đủ hai tệp đọc được",
                "bo_qua": bo_qua}

    try:
        if dinh_dang == "md":
            khoi = []
            for q in hop_le:
                khoi.append(f"# {q.stem}\n\n{_doc_chu(q)}\n")
            ra.write_text("\n---\n\n".join(khoi), "utf-8")
        elif dinh_dang == "docx":
            _noi_docx(hop_le, ra)
        else:
            return {"ok": False, "error": f"chỉ nối được docx hoặc md, không phải {dinh_dang}"}
    except Exception as exc:
        logger.warning("noi_tep lỗi: %s", str(exc)[:150])
        return {"ok": False, "error": f"không nối được: {str(exc)[:150]}"}
    kq: dict[str, Any] = {"ok": True, "duong_dan": str(ra), "ten": ra.name,
                          "so_tep": len(hop_le)}
    if bo_qua:
        # Nói RA tệp nào đã rơi. Bỏ qua im lặng là nối 5 tệp ra 2 mà không ai biết.
        kq["bo_qua"] = bo_qua
    return kq


def _noi_docx(cac_tep: list[Path], ra: Path) -> None:
    """Nối .docx GIỮ ĐỊNH DẠNG: bê nguyên thân XML của từng đoạn/bảng sang.

    Repo gán `new_para.text = ...` nên mất đậm/nghiêng và bỏ hết bảng. Ở đây
    dùng `copy.deepcopy` phần tử XML rồi chèn vào thân tài liệu đích — giữ được
    cả định dạng chữ lẫn bảng.
    """
    import copy
    import docx
    ra_doc = docx.Document()
    than = ra_doc.element.body
    for i, q in enumerate(cac_tep):
        ra_doc.add_heading(q.stem, level=1)
        if q.suffix.lower() == ".docx":
            nguon = docx.Document(str(q))
            for pt in nguon.element.body:
                if pt.tag.endswith("}sectPr"):     # thẻ khổ giấy — không bê sang
                    continue
                than.append(copy.deepcopy(pt))
        else:
            for d in _doc_chu(q).splitlines():
                ra_doc.add_paragraph(d)
        if i < len(cac_tep) - 1:
            ra_doc.add_page_break()
    ra_doc.save(str(ra))


# ── 7. Đọc thông tin tài liệu ──────────────────────────────────────────────

def doc_thong_tin(tep: str) -> dict[str, Any]:
    """Thông tin một tài liệu: cỡ, lần sửa, tác giả, số trang/sheet/slide/từ."""
    try:
        p = _duong_dan(tep)
    except (ValueError, FileNotFoundError, OSError) as exc:
        return {"ok": False, "error": str(exc)}
    from datetime import datetime
    st = p.stat()
    tt: dict[str, Any] = {
        "ten": p.name, "loai": p.suffix.lower().lstrip("."),
        "co_byte": st.st_size,
        "sua_luc": datetime.fromtimestamp(st.st_mtime).strftime("%d/%m/%Y %H:%M"),
    }
    duoi = p.suffix.lower()
    try:
        if duoi == ".docx":
            import docx
            d = docx.Document(str(p))
            cp = d.core_properties
            tt.update({"tac_gia": cp.author or "", "tieu_de": cp.title or "",
                       "so_doan": len(d.paragraphs), "so_bang": len(d.tables),
                       "so_tu": sum(len(x.text.split()) for x in d.paragraphs)})
            if cp.created:
                tt["tao_luc"] = cp.created.strftime("%d/%m/%Y %H:%M")
        elif duoi in {".xlsx", ".xlsm"}:
            import openpyxl
            wb = openpyxl.load_workbook(str(p), read_only=True)
            tt.update({"so_sheet": len(wb.worksheets),
                       "cac_sheet": [w.title for w in wb.worksheets],
                       "so_dong": {w.title: (w.max_row or 0) for w in wb.worksheets}})
            wb.close()
        elif duoi == ".pptx":
            from pptx import Presentation
            pres = Presentation(str(p))
            tt["so_slide"] = len(pres.slides)
        elif duoi == ".pdf":
            import fitz
            with fitz.open(str(p)) as d:
                md = d.metadata or {}
                tt.update({"so_trang": d.page_count,
                           "tac_gia": md.get("author") or "",
                           "tieu_de": md.get("title") or ""})
    except Exception as exc:
        # Đọc được cỡ/ngày rồi thì vẫn trả về, chỉ ghi chú phần đọc sâu bị hỏng.
        tt["ghi_chu"] = f"không đọc được phần chi tiết: {str(exc)[:100]}"
    return {"ok": True, **tt}


# ── 8. Sinh báo cáo từ Excel + mẫu Markdown ────────────────────────────────
# Repo nhận tham số `template_md` rồi KHÔNG DÙNG ở đâu cả (soi hết
# report_creator.py: chỉ xuất hiện ở chữ ký hàm và một lời gọi chuyển tiếp), nên
# tiêu đề báo cáo luôn cứng là "Báo cáo: <tên tệp>". Ở đây mẫu là thật: mỗi dòng
# Markdown thành một đoạn, và chỗ nào ghi `{{bang:<tên sheet>}}` thì thay bằng
# bảng thật; `{{thong_ke:<tên sheet>}}` thì thay bằng bảng min/max/trung bình.

_RE_CHO_THAY = re.compile(r"\{\{\s*(bang|thong_ke)\s*:\s*([^}]+?)\s*\}\}")


def tao_bao_cao(tep_du_lieu: str, ten_ra: str, *,
                mau_md: str = "", max_dong_bang: int = 50) -> dict[str, Any]:
    """Excel → báo cáo .docx. `mau_md` là mẫu Markdown, để trống thì tự dựng."""
    try:
        p = _duong_dan(tep_du_lieu)
        from services import officecli
        ra = officecli.resolve_path(str(ten_ra or "").strip() or "bao-cao.docx")
    except (ValueError, FileNotFoundError, OSError) as exc:
        return {"ok": False, "error": str(exc)}
    if p.suffix.lower() not in {".xlsx", ".xls", ".xlsm", ".csv"}:
        return {"ok": False, "error": f"cần bảng tính, không phải {p.suffix}"}
    try:
        import pandas as pd
        if p.suffix.lower() == ".csv":
            bang = {"CSV": pd.read_csv(str(p))}
        else:
            bang = pd.read_excel(str(p), sheet_name=None)
    except Exception as exc:
        return {"ok": False, "error": f"không đọc được bảng: {str(exc)[:150]}"}
    if not bang:
        return {"ok": False, "error": "bảng tính không có sheet nào"}

    try:
        import docx
        d = docx.Document()
        if str(mau_md or "").strip():
            _theo_mau(d, mau_md, bang, max_dong_bang)
        else:
            _bao_cao_mac_dinh(d, p.stem, bang, max_dong_bang)
        d.save(str(ra))
    except Exception as exc:
        logger.warning("tao_bao_cao lỗi: %s", str(exc)[:150])
        return {"ok": False, "error": f"không tạo được báo cáo: {str(exc)[:150]}"}
    return {"ok": True, "duong_dan": str(ra), "ten": ra.name,
            "so_sheet": len(bang), "theo_mau": bool(str(mau_md or "").strip())}


def _bang_vao_docx(d: Any, df: Any, max_dong: int) -> None:
    cot = [str(c) for c in df.columns]
    b = d.add_table(rows=1, cols=len(cot))
    b.style = "Light Grid Accent 1"
    for i, c in enumerate(cot):
        b.rows[0].cells[i].text = c
    for _, hang in df.head(max_dong).iterrows():
        o = b.add_row().cells
        for i, c in enumerate(cot):
            v = hang[c]
            o[i].text = "" if v is None or (isinstance(v, float) and v != v) else str(v)
    if len(df) > max_dong:
        d.add_paragraph(f"(còn {len(df) - max_dong} dòng nữa, không đưa vào báo cáo)")


def _thong_ke_vao_docx(d: Any, df: Any) -> None:
    so = df.select_dtypes("number")
    if so.empty:
        d.add_paragraph("(sheet này không có cột số nào để thống kê)")
        return
    b = d.add_table(rows=1, cols=5)
    b.style = "Light Grid Accent 1"
    for i, t in enumerate(["Cột", "Nhỏ nhất", "Lớn nhất", "Trung bình", "Tổng"]):
        b.rows[0].cells[i].text = t
    for c in so.columns:
        v = so[c].dropna()
        if v.empty:
            continue
        o = b.add_row().cells
        for i, t in enumerate([str(c), f"{v.min():g}", f"{v.max():g}",
                               f"{v.mean():.2f}", f"{v.sum():g}"]):
            o[i].text = t


def _theo_mau(d: Any, mau: str, bang: dict, max_dong: int) -> None:
    for dong in str(mau).splitlines():
        m = _RE_CHO_THAY.search(dong)
        if m:
            loai, ten_sheet = m.group(1), m.group(2).strip()
            df = bang.get(ten_sheet)
            if df is None:
                d.add_paragraph(f"[không có sheet “{ten_sheet}” trong tệp]")
            elif loai == "bang":
                _bang_vao_docx(d, df, max_dong)
            else:
                _thong_ke_vao_docx(d, df)
            continue
        td = _RE_TD.match(dong.strip())
        if td:
            d.add_heading(td.group(2), level=min(len(td.group(1)), 4))
        elif dong.strip():
            g = _RE_GACH_DAU.match(dong)
            d.add_paragraph(g.group(1) if g else dong.strip(),
                            style="List Bullet" if g else None)


def _bao_cao_mac_dinh(d: Any, ten: str, bang: dict, max_dong: int) -> None:
    from datetime import datetime
    d.add_heading(f"Báo cáo {ten}", level=1)
    d.add_paragraph("Lập lúc " + datetime.now().strftime("%H:%M %d/%m/%Y"))
    for ten_sheet, df in bang.items():
        d.add_heading(str(ten_sheet), level=2)
        if df is None or df.empty:
            d.add_paragraph("(trống)")
            continue
        d.add_paragraph(f"{len(df)} dòng × {len(df.columns)} cột")
        d.add_heading("Thống kê", level=3)
        _thong_ke_vao_docx(d, df)
        d.add_heading("Dữ liệu", level=3)
        _bang_vao_docx(d, df, max_dong)
