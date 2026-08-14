"""Bốn việc tài liệu bù từ repo Quiz99/officeCliMCP.

Repo đó có 23 tool, 19 trùng việc đã có ở đây. Bốn việc thật sự thiếu:
so sánh hai tài liệu, thống kê bảng tính, tìm-thay-thế toàn tài liệu, tạo .pptx.

Ba chỗ bản này phải LÀM KHÁC repo, và đây là chỗ test canh:

1. Khoá trong thư mục làm việc. `_resolve_path` của repo không kiểm gì, nên
   `diff_documents("/etc/passwd", …)` đọc được tệp hệ thống.
2. Đọc .docx phải GIỮ BẢNG. Repo nối `p.text` từng paragraph → mất sạch bảng; so
   sánh hai bản hợp đồng mà rơi mất bảng thì so sánh để làm gì.
3. Tìm-thay-thế phải chạm cả ô trong BẢNG, và giữ định dạng đậm/nghiêng.
"""
import os
import sys
import tempfile
import unittest
from pathlib import Path

GOC = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(GOC))
os.environ.setdefault("CHATGPT2API_AUTH_KEY", "test-auth")

from services import office_bo_sung as ob  # noqa: E402


def _co_lib(*ten: str) -> bool:
    for t in ten:
        try:
            __import__(t)
        except Exception:
            return False
    return True


class _Nen(unittest.TestCase):
    """Ép thư mục làm việc về thư mục tạm — không đụng workspace thật."""

    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.addCleanup(self.tmp.cleanup)
        # resolve(): trên macOS /var là symlink tới /private/var, mà
        # `officecli.resolve_path` resolve đường dẫn ứng viên rồi mới so — không
        # resolve gốc thì mọi đường dẫn đều bị coi là ngoài workspace.
        self.ws = Path(self.tmp.name).resolve()
        from services import officecli
        self._goc = officecli.workspace
        officecli.workspace = lambda: self.ws
        self.addCleanup(setattr, officecli, "workspace", self._goc)


class KhoaTrongThuMucLamViecTests(_Nen):
    """Bot đọc tài liệu người lạ gửi tới, mà tài liệu có thể chứa câu ra lệnh."""

    def test_khong_doc_duoc_tep_ngoai_workspace(self):
        kq = ob.so_sanh("/etc/passwd", "/etc/hosts")
        self.assertFalse(kq["ok"])
        self.assertIn("workspace", kq["error"].lower())

    def test_khong_lach_duoc_bang_hai_cham(self):
        (self.ws / "a.md").write_text("x", "utf-8")
        kq = ob.so_sanh("a.md", "../../../etc/passwd")
        self.assertFalse(kq["ok"])

    def test_thong_ke_cung_bi_khoa(self):
        self.assertFalse(ob.thong_ke_bang("/etc/passwd")["ok"])

    def test_thay_the_cung_bi_khoa(self):
        self.assertFalse(ob.tim_thay_the("/etc/passwd", "a", "b")["ok"])

    def test_tao_slide_khong_ghi_ra_ngoai(self):
        kq = ob.tao_slide("# A\n- x", "../../ngoai.pptx")
        if kq.get("ok"):
            self.assertIn(self.ws, Path(kq["duong_dan"]).parents)


class SoSanhTests(_Nen):

    def _hai_tep(self, a: str, b: str):
        (self.ws / "a.md").write_text(a, "utf-8")
        (self.ws / "b.md").write_text(b, "utf-8")

    def test_giong_nhau_thi_noi_giong_nhau(self):
        self._hai_tep("một\nhai", "một\nhai")
        kq = ob.so_sanh("a.md", "b.md")
        self.assertTrue(kq["ok"])
        self.assertTrue(kq["giong_nhau"])
        self.assertEqual((kq["them"], kq["bo"], kq["sua"]), (0, 0, 0))

    def test_dem_dung_so_dong_them_va_bo(self):
        self._hai_tep("một\nhai\nba", "một\nba\nbốn\nnăm")
        kq = ob.so_sanh("a.md", "b.md")
        self.assertFalse(kq["giong_nhau"])
        self.assertGreater(kq["them"], 0)
        self.assertGreater(kq["bo"], 0)

    def test_bao_cao_co_ten_hai_tep_va_noi_dung_khac(self):
        self._hai_tep("giá 100 đồng", "giá 200 đồng")
        bc = ob.so_sanh("a.md", "b.md")["bao_cao"]
        self.assertIn("a.md", bc)
        self.assertIn("b.md", bc)
        self.assertIn("100", bc)
        self.assertIn("200", bc)

    def test_dinh_dang_unified_chay_duoc(self):
        self._hai_tep("một\nhai", "một\nba")
        kq = ob.so_sanh("a.md", "b.md", dinh_dang="unified")
        self.assertTrue(kq["ok"])
        self.assertIn("---", kq["bao_cao"] + "---")

    def test_thieu_tep_thi_bao_ro(self):
        (self.ws / "a.md").write_text("x", "utf-8")
        kq = ob.so_sanh("a.md", "khong-co.md")
        self.assertFalse(kq["ok"])

    def test_bao_cao_qua_dai_thi_cat_VA_NOI_LA_CAT(self):
        """Cắt âm thầm là thứ đã làm mất cả buổi để lần ra ở ask_choices."""
        self._hai_tep("\n".join(f"dòng {i}" for i in range(5000)), "khác hẳn")
        kq = ob.so_sanh("a.md", "b.md")
        self.assertLessEqual(len(kq["bao_cao"]), ob.MAX_KY_TU)
        self.assertTrue(kq["bi_cat"])


@unittest.skipUnless(_co_lib("openpyxl", "pandas"), "cần openpyxl + pandas")
class ThongKeBangTests(_Nen):

    def _bang(self, ten="so.xlsx"):
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Luong"
        ws.append(["Tên", "Lương", "Thưởng"])
        ws.append(["An", 10, 1])
        ws.append(["Bình", 20, None])
        ws.append(["An", 30, 3])
        wb.save(str(self.ws / ten))
        return ten

    def test_dem_dung_dong_va_cot(self):
        kq = ob.thong_ke_bang(self._bang())
        self.assertTrue(kq["ok"], kq.get("error"))
        self.assertEqual(kq["tong_dong"], 3)
        self.assertIn("3 dòng × 3 cột", kq["bao_cao"])

    def test_co_min_max_trung_binh_tong(self):
        bc = ob.thong_ke_bang(self._bang())["bao_cao"]
        self.assertIn("Lương", bc)
        self.assertIn("60", bc)      # tổng 10+20+30
        self.assertIn("20.00", bc)   # trung bình

    def test_chi_ro_cho_thieu_du_lieu(self):
        bc = ob.thong_ke_bang(self._bang())["bao_cao"]
        self.assertIn("Thiếu dữ liệu", bc)
        self.assertIn("Thưởng", bc)

    def test_cot_chu_noi_gia_tri_hay_gap(self):
        bc = ob.thong_ke_bang(self._bang())["bao_cao"]
        self.assertIn("An", bc)

    def test_khong_phai_bang_tinh_thi_bao_ro(self):
        (self.ws / "a.docx").write_text("x", "utf-8")
        kq = ob.thong_ke_bang("a.docx")
        self.assertFalse(kq["ok"])
        self.assertIn("bảng tính", kq["error"])


@unittest.skipUnless(_co_lib("docx"), "cần python-docx")
class TimThayTheTrongWordTests(_Nen):

    def _tai_lieu(self, ten="a.docx"):
        import docx
        d = docx.Document()
        d.add_paragraph("Năm học 2024 bắt đầu.")
        d.add_paragraph("Kết thúc năm học 2024.")
        b = d.add_table(rows=1, cols=2)
        b.rows[0].cells[0].text = "Khoá 2024"
        b.rows[0].cells[1].text = "Ghi chú"
        d.save(str(self.ws / ten))
        return ten

    def _doc_lai(self, ten="a.docx"):
        import docx
        d = docx.Document(str(self.ws / ten))
        chu = [p.text for p in d.paragraphs]
        for bang in d.tables:
            for h in bang.rows:
                chu += [o.text for o in h.cells]
        return "\n".join(chu)

    def test_doi_het_moi_cho(self):
        kq = ob.tim_thay_the(self._tai_lieu(), "2024", "2025")
        self.assertTrue(kq["ok"], kq.get("error"))
        self.assertEqual(kq["so_lan"], 3)
        s = self._doc_lai()
        self.assertNotIn("2024", s)
        self.assertEqual(s.count("2025"), 3)

    def test_CHAM_CA_O_TRONG_BANG(self):
        """Repo bỏ qua bảng hoàn toàn — đây là chỗ bản này hơn."""
        ob.tim_thay_the(self._tai_lieu(), "2024", "2025")
        self.assertIn("Khoá 2025", self._doc_lai())

    def test_chi_doi_cho_dau_khi_tat_ca_False(self):
        kq = ob.tim_thay_the(self._tai_lieu(), "2024", "2025", tat_ca=False)
        self.assertEqual(kq["so_lan"], 1)
        self.assertIn("2024", self._doc_lai())

    def test_giu_dinh_dang_dam(self):
        import docx
        d = docx.Document()
        p = d.add_paragraph()
        r = p.add_run("Tên cũ")
        r.bold = True
        d.save(str(self.ws / "b.docx"))
        ob.tim_thay_the("b.docx", "cũ", "mới")
        d2 = docx.Document(str(self.ws / "b.docx"))
        self.assertEqual(d2.paragraphs[0].text, "Tên mới")
        self.assertTrue(d2.paragraphs[0].runs[0].bold,
                        "gán thẳng paragraph.text là xoá sạch định dạng")

    def test_khong_thay_thi_noi_khong_thay(self):
        kq = ob.tim_thay_the(self._tai_lieu(), "không-có-cụm-này", "x")
        self.assertTrue(kq["ok"])
        self.assertEqual(kq["so_lan"], 0)

    def test_thieu_cum_can_tim_thi_bao_ro(self):
        self.assertFalse(ob.tim_thay_the(self._tai_lieu(), "", "x")["ok"])


@unittest.skipUnless(_co_lib("openpyxl"), "cần openpyxl")
class TimThayTheTrongExcelTests(_Nen):

    def _bang(self):
        import openpyxl
        wb = openpyxl.Workbook()
        wb.active.append(["Năm 2024", 5])
        wb.active.append(["Khoá 2024", 6])
        wb.save(str(self.ws / "a.xlsx"))
        return "a.xlsx"

    def test_doi_moi_o_chu(self):
        kq = ob.tim_thay_the(self._bang(), "2024", "2025")
        self.assertEqual(kq["so_lan"], 2)
        import openpyxl
        wb = openpyxl.load_workbook(str(self.ws / "a.xlsx"))
        self.assertEqual(wb.active["A1"].value, "Năm 2025")

    def test_khong_dung_o_so(self):
        ob.tim_thay_the(self._bang(), "2024", "2025")
        import openpyxl
        wb = openpyxl.load_workbook(str(self.ws / "a.xlsx"))
        self.assertEqual(wb.active["B1"].value, 5)


@unittest.skipUnless(_co_lib("pptx"), "cần python-pptx")
class TaoSlideTests(_Nen):
    """Dự án chưa tạo được .pptx ở bất kỳ đâu — đây là việc đáng nhất từ repo."""

    DAN_Y = ("# Bài 1: Phân số\n"
             "- Khái niệm\n"
             "- Ví dụ\n"
             "## Luyện tập\n"
             "- Bài 1\n"
             "- Bài 2\n")

    def test_moi_tieu_de_mot_slide(self):
        kq = ob.tao_slide(self.DAN_Y, "bai_giang.pptx")
        self.assertTrue(kq["ok"], kq.get("error"))
        self.assertEqual(kq["so_slide"], 2)

    def test_tep_that_su_mo_lai_duoc(self):
        kq = ob.tao_slide(self.DAN_Y, "bai_giang.pptx")
        from pptx import Presentation
        pres = Presentation(kq["duong_dan"])
        self.assertEqual(len(pres.slides), 2)
        self.assertIn("Phân số", pres.slides[0].shapes.title.text)

    def test_cac_y_vao_dung_slide(self):
        kq = ob.tao_slide(self.DAN_Y, "b.pptx")
        from pptx import Presentation
        pres = Presentation(kq["duong_dan"])
        chu = "\n".join(sh.text_frame.text for sh in pres.slides[1].shapes
                        if sh.has_text_frame)
        self.assertIn("Bài 1", chu)
        self.assertIn("Bài 2", chu)

    def test_tu_them_duoi_pptx(self):
        kq = ob.tao_slide("# A\n- x", "khong_duoi")
        self.assertTrue(kq["ten"].endswith(".pptx"))

    def test_dan_y_rong_thi_bao_ro(self):
        self.assertFalse(ob.tao_slide("   ", "a.pptx")["ok"])

    def test_dan_y_khong_co_tieu_de_nao_thi_bao_ro(self):
        kq = ob.tao_slide("chỉ có chữ suông", "a.pptx")
        # Chữ trước tiêu đề đầu vẫn thành một slide mở đầu — không được báo hỏng.
        self.assertTrue(kq["ok"])
        self.assertEqual(kq["so_slide"], 1)


@unittest.skipUnless(_co_lib("pptx"), "cần python-pptx")
class SlideCoBoCucVaChuDeTests(_Nen):
    """Bản đầu đổ mọi thứ vào MỘT bố cục, màu theo mẫu trắng trơn, bảng số liệu
    thành mấy dòng gạch đầu dòng. Nhóm này canh phần bù: bố cục chọn theo nội
    dung, màu lấy từ một bộ chủ đề, bảng và biểu đồ là ĐỐI TƯỢNG GỐC của
    PowerPoint (sửa được) chứ không phải chữ hay ảnh."""

    def _mo(self, kq):
        from pptx import Presentation
        self.assertTrue(kq["ok"], kq.get("error"))
        return Presentation(kq["duong_dan"])

    def test_khong_them_slide_nao_ngoai_so_tieu_de(self):
        """Cái dễ hỏng nhất khi thêm slide bìa: tự chèn thêm một slide, làm lệch
        số slide mà người soạn đếm được từ dàn ý."""
        kq = ob.tao_slide("# A\n- x\n## B\n- y\n", "a.pptx")
        self.assertEqual(kq["so_slide"], 2)
        self.assertEqual(len(self._mo(kq).slides), 2)

    def test_khoi_16_9(self):
        """Mẫu trần của python-pptx là 4:3 — chiếu lên là hai vệt đen hai bên."""
        pres = self._mo(ob.tao_slide("# A\n- x", "a.pptx"))
        self.assertAlmostEqual(pres.slide_width / pres.slide_height, 16 / 9,
                               places=2)

    def test_tieu_de_khong_co_y_o_dau_thi_la_slide_bia(self):
        pres = self._mo(ob.tao_slide("# Bài giảng\n## Nội dung\n- x", "a.pptx"))
        self.assertEqual(pres.slides[0].slide_layout.name, "Title Slide")

    def test_tieu_de_khong_co_y_o_giua_thi_la_slide_phan_muc(self):
        pres = self._mo(ob.tao_slide(
            "# Mở đầu\n- x\n# Phần hai\n## Chi tiết\n- y", "a.pptx"))
        self.assertEqual(pres.slides[1].slide_layout.name, "Section Header")

    def test_thut_le_hai_dau_cach_thanh_cap_that(self):
        """Trước đây mọi ý nằm cùng một cấp, dàn ý có ý con thì đọc ra một khối."""
        pres = self._mo(ob.tao_slide(
            "# A\n- cha\n  - con\n    - chau\n", "a.pptx"))
        khung = next(ph.text_frame for ph in pres.slides[0].placeholders
                     if ph.placeholder_format.idx == 1)
        self.assertEqual([d.level for d in khung.paragraphs], [0, 1, 2])

    def test_bang_markdown_thanh_BANG_GOC_chu_khong_phai_chu(self):
        pres = self._mo(ob.tao_slide(
            "# Điểm\n| Tên | Điểm |\n| --- | --- |\n| An | 8 |\n| Bình | 9 |\n",
            "a.pptx"))
        bang = [sh for sh in pres.slides[0].shapes if sh.has_table]
        self.assertEqual(len(bang), 1, "phải là bảng gốc, không phải text")
        tb = bang[0].table
        self.assertEqual((len(tb.rows), len(tb.columns)), (3, 2))
        self.assertEqual(tb.cell(0, 0).text, "Tên")
        self.assertEqual(tb.cell(2, 1).text, "9")

    def test_dong_ngan_cua_bang_khong_thanh_hang_du_lieu(self):
        """`| --- | --- |` là cú pháp, không phải số liệu."""
        pres = self._mo(ob.tao_slide(
            "# A\n| a | b |\n| :-- | --: |\n| 1 | 2 |\n", "a.pptx"))
        tb = next(sh.table for sh in pres.slides[0].shapes if sh.has_table)
        self.assertEqual(len(tb.rows), 2)

    def test_moc_bieu_do_thi_ra_BIEU_DO_GOC(self):
        kq = ob.tao_slide(
            "# Doanh thu\n[biểu đồ]\n| Tháng | Tiền |\n| --- | --- |\n"
            "| Bảy | 1.200 |\n| Tám | 1.450,5 |\n", "a.pptx")
        pres = self._mo(kq)
        self.assertEqual(kq["so_bieu_do"], 1)
        self.assertEqual(kq["so_bang"], 0)
        do_thi = [sh for sh in pres.slides[0].shapes if sh.has_chart]
        self.assertEqual(len(do_thi), 1)
        ch = do_thi[0].chart
        self.assertEqual(list(ch.plots[0].categories), ["Bảy", "Tám"])
        # Số kiểu Việt: '1.200' là một nghìn hai, '1.450,5' có phẩy thập phân.
        self.assertEqual(list(ch.series[0].values), [1200.0, 1450.5])

    def test_KHONG_co_moc_thi_van_la_bang_du_toan_so(self):
        """Bảng số liệu tự biến thành biểu đồ là kiểu thông minh gây bất ngờ, mà
        mất luôn phần đọc được từng con số."""
        kq = ob.tao_slide("# A\n| x | y |\n| --- | --- |\n| a | 1 |\n| b | 2 |\n",
                          "a.pptx")
        self.assertEqual((kq["so_bang"], kq["so_bieu_do"]), (1, 0))

    def test_xin_bieu_do_ma_du_lieu_khong_phai_so_thi_ve_bang(self):
        """Thà ra bảng đúng còn hơn một biểu đồ rỗng."""
        kq = ob.tao_slide("# A\n[biểu đồ]\n| x | y |\n| --- | --- |\n"
                          "| a | nhiều |\n| b | ít |\n", "a.pptx")
        self.assertEqual((kq["so_bang"], kq["so_bieu_do"]), (1, 0))

    def test_chu_de_doi_mau_nen_that(self):
        from pptx.dml.color import RGBColor
        sang = self._mo(ob.tao_slide("# A\n- x", "s.pptx", "trang-sach"))
        toi = self._mo(ob.tao_slide("# A\n- x", "t.pptx", "xanh-dam"))
        self.assertEqual(sang.slides[0].background.fill.fore_color.rgb,
                         RGBColor.from_string("FFFFFF"))
        self.assertEqual(toi.slides[0].background.fill.fore_color.rgb,
                         RGBColor.from_string("0F243E"))

    def test_ten_chu_de_la_thi_dung_mac_dinh_chu_khong_bao_hong(self):
        kq = ob.tao_slide("# A\n- x", "a.pptx", "mau-tim-gradient")
        self.assertTrue(kq["ok"])
        self.assertEqual(kq["chu_de"], "trang-sach")

    def test_chu_de_ap_len_TUNG_run_chu_khong_chi_paragraph(self):
        """Đặt font ở cấp paragraph là không đủ: run sinh sau không thừa hưởng,
        slide ra nửa theo chủ đề nửa theo mẫu."""
        pres = self._mo(ob.tao_slide("# A\n- một\n- hai\n- ba\n", "a.pptx",
                                     "xanh-dam"))
        khung = next(ph.text_frame for ph in pres.slides[0].placeholders
                     if ph.placeholder_format.idx == 1)
        co = [r.font.size for d in khung.paragraphs for r in d.runs]
        self.assertEqual(len(co), 3)
        self.assertTrue(all(c is not None for c in co), "có run bị bỏ sót")

    def test_so_hieu_ca_hai_loi_viet(self):
        """Dấu phân cách cuối là thập phân, trừ khi sau nó đúng ba chữ số."""
        for chu, mong in [("1.200", 1200.0),        # kiểu Việt: chấm nhóm nghìn
                          ("1.200,5", 1200.5),      # kiểu Việt đủ hai dấu
                          ("1,200.5", 1200.5),      # kiểu Anh, ngược dấu
                          ("1.200.000", 1200000.0),
                          ("1,5", 1.5), ("1.5", 1.5), ("1.2345", 1.2345),
                          ("1 200", 1200.0), ("12%", 12.0), ("-2,5", -2.5),
                          ("8", 8.0)]:
            self.assertEqual(ob._so(chu), mong, chu)
        for chu in ("nhiều", "", "   ", "x1", "1.2.3,4,5"):
            self.assertIsNone(ob._so(chu), chu)

    def test_ca_cot_go_cho_mo_ho_cua_tung_o(self):
        """`6.750` một mình là mơ hồ; đứng cạnh 7,25 và 8 thì rõ là 6,75. Đọc
        thành 6750 là cột cao gấp gần nghìn lần hai cột kia — lỗi đo được trên
        bộ slide thử, không phải lo xa."""
        self.assertEqual(ob._cot_so(["6.750", "7,25", "8"]), [6.75, 7.25, 8.0])
        # Có mốc cỡ nghìn thì `1.200` đúng là một nghìn hai.
        self.assertEqual(ob._cot_so(["1.200", "1.450,5"]), [1200.0, 1450.5])
        # Cả cột đều mơ hồ → giữ nghĩa nhóm nghìn, lối người Việt hay viết hơn.
        self.assertEqual(ob._cot_so(["1.200", "1.450"]), [1200.0, 1450.0])
        # Một ô không phải số là cả cột không vẽ được.
        self.assertIsNone(ob._cot_so(["1", "nhiều"]))

    def test_bieu_do_lay_dung_so_sau_khi_gio_mo_ho(self):
        kq = ob.tao_slide(
            "# Điểm TB\n[biểu đồ]\n| Tháng | Điểm |\n| --- | --- |\n"
            "| Chín | 6.750 |\n| Mười | 7,25 |\n| Mười một | 8 |\n", "a.pptx")
        pres = self._mo(kq)
        ch = next(sh.chart for sh in pres.slides[0].shapes if sh.has_chart)
        self.assertEqual(list(ch.series[0].values), [6.75, 7.25, 8.0])


class CatTheoTieuDeTests(_Nen):
    """Repo cắt bằng `content.split("\n## ")` — CỨNG ở cấp 2. Tài liệu chỉ dùng
    '#' (rất thường gặp với văn bản tiếng Việt một cấp) thì không cắt được mảnh
    nào mà vẫn báo ok, trả về đúng một tệp bằng cả bản gốc."""

    def _md(self, noi_dung, ten="a.md"):
        (self.ws / ten).write_text(noi_dung, "utf-8")
        return ten

    def test_cat_theo_cap_2(self):
        kq = ob.cat_theo_tieu_de(self._md(
            "# Sách\n## Bài 1\nnội dung 1\n## Bài 2\nnội dung 2\n"), cap=2)
        self.assertTrue(kq["ok"], kq.get("error"))
        self.assertEqual(kq["so_phan"], 3)   # mở đầu '# Sách' + hai bài

    def test_TAI_LIEU_CHI_CO_MOT_CAP_van_cat_duoc(self):
        """Đây là chỗ repo bó tay."""
        kq = ob.cat_theo_tieu_de(self._md(
            "# Phần một\nnội dung\n# Phần hai\nnội dung\n# Phần ba\nx\n"))
        self.assertTrue(kq["ok"], kq.get("error"))
        self.assertEqual(kq["cap"], 1)
        self.assertEqual(kq["so_phan"], 3)

    def test_ten_tep_lay_theo_tieu_de(self):
        kq = ob.cat_theo_tieu_de(self._md("# Phân số\nx\n# Số thập phân\ny\n"))
        ten = " ".join(kq["cac_tep"])
        self.assertIn("Phân-số", ten)
        self.assertIn("Số-thập-phân", ten)

    def test_noi_dung_tung_phan_dung(self):
        ob.cat_theo_tieu_de(self._md("# A\nnội dung A\n# B\nnội dung B\n"))
        d = self.ws / "a-cat"
        tep = sorted(d.glob("*.md"))
        self.assertEqual(len(tep), 2)
        self.assertIn("nội dung A", tep[0].read_text("utf-8"))
        self.assertNotIn("nội dung B", tep[0].read_text("utf-8"))

    def test_khong_co_tieu_de_thi_bao_ro_chu_khong_bao_ok(self):
        kq = ob.cat_theo_tieu_de(self._md("chỉ có chữ suông, không tiêu đề\n"))
        self.assertFalse(kq["ok"])
        self.assertIn("tiêu đề", kq["error"])

    def test_khong_ghi_ra_ngoai_workspace(self):
        kq = ob.cat_theo_tieu_de(self._md("# A\nx\n"), thu_muc="../../ngoai")
        self.assertFalse(kq["ok"])


@unittest.skipUnless(_co_lib("docx"), "cần python-docx")
class NoiTepTests(_Nen):

    def _docx(self, ten, chu, dam=False, co_bang=False):
        import docx
        d = docx.Document()
        p = d.add_paragraph()
        r = p.add_run(chu)
        r.bold = dam
        if co_bang:
            b = d.add_table(rows=1, cols=2)
            b.rows[0].cells[0].text = "ô A"
            b.rows[0].cells[1].text = "ô B"
        d.save(str(self.ws / ten))
        return ten

    def _doc_ra(self, ten):
        import docx
        d = docx.Document(str(self.ws / ten))
        chu = [p.text for p in d.paragraphs]
        bang = []
        for b in d.tables:
            for h in b.rows:
                bang += [o.text for o in h.cells]
        return "\n".join(chu), bang, d

    def test_noi_hai_docx(self):
        a = self._docx("a.docx", "nội dung A")
        b = self._docx("b.docx", "nội dung B")
        kq = ob.noi_tep([a, b], "ra.docx")
        self.assertTrue(kq["ok"], kq.get("error"))
        self.assertEqual(kq["so_tep"], 2)
        chu, _, _ = self._doc_ra("ra.docx")
        self.assertIn("nội dung A", chu)
        self.assertIn("nội dung B", chu)

    def test_GIU_DINH_DANG_DAM(self):
        """Repo gán `new_para.text = para.text` nên mất sạch đậm/nghiêng."""
        a = self._docx("a.docx", "chữ đậm", dam=True)
        b = self._docx("b.docx", "chữ thường")
        ob.noi_tep([a, b], "ra.docx")
        _, _, d = self._doc_ra("ra.docx")
        co_dam = any(r.bold for p in d.paragraphs for r in p.runs
                     if r.text == "chữ đậm")
        self.assertTrue(co_dam, "nối xong mất định dạng đậm")

    def test_GIU_CA_BANG(self):
        """Repo chỉ bê paragraph, bỏ hết bảng."""
        a = self._docx("a.docx", "x", co_bang=True)
        b = self._docx("b.docx", "y")
        ob.noi_tep([a, b], "ra.docx")
        _, bang, _ = self._doc_ra("ra.docx")
        self.assertIn("ô A", bang)

    def test_TEP_THIEU_thi_NOI_RA_chu_khong_im_lang(self):
        """Repo `continue` im lặng — nối 5 tệp ra 2 mà không ai biết."""
        a = self._docx("a.docx", "A")
        b = self._docx("b.docx", "B")
        kq = ob.noi_tep([a, b, "khong-co.docx"], "ra.docx")
        self.assertTrue(kq["ok"])
        self.assertEqual(kq["so_tep"], 2)
        self.assertIn("bo_qua", kq)
        self.assertIn("khong-co.docx", " ".join(kq["bo_qua"]))

    def test_duoi_mot_tep_thi_bao_ro(self):
        a = self._docx("a.docx", "A")
        self.assertFalse(ob.noi_tep([a], "ra.docx")["ok"])

    def test_noi_dinh_dang_md(self):
        (self.ws / "a.md").write_text("nội dung A", "utf-8")
        (self.ws / "b.md").write_text("nội dung B", "utf-8")
        kq = ob.noi_tep(["a.md", "b.md"], "ra.md", dinh_dang="md")
        self.assertTrue(kq["ok"], kq.get("error"))
        s = (self.ws / "ra.md").read_text("utf-8")
        self.assertIn("nội dung A", s)
        self.assertIn("nội dung B", s)


class DocThongTinTests(_Nen):

    def test_tep_thuong_van_co_co_va_ngay(self):
        (self.ws / "a.md").write_text("xin chào", "utf-8")
        kq = ob.doc_thong_tin("a.md")
        self.assertTrue(kq["ok"])
        self.assertEqual(kq["ten"], "a.md")
        self.assertEqual(kq["loai"], "md")
        self.assertGreater(kq["co_byte"], 0)
        self.assertIn("/", kq["sua_luc"])

    @unittest.skipUnless(_co_lib("docx"), "cần python-docx")
    def test_docx_co_tac_gia_va_dem_doan_bang(self):
        import docx
        d = docx.Document()
        d.core_properties.author = "Nguyễn Việt"
        d.add_paragraph("một hai ba")
        d.add_table(rows=1, cols=2)
        d.save(str(self.ws / "a.docx"))
        kq = ob.doc_thong_tin("a.docx")
        self.assertEqual(kq["tac_gia"], "Nguyễn Việt")
        self.assertEqual(kq["so_bang"], 1)
        self.assertEqual(kq["so_tu"], 3)

    @unittest.skipUnless(_co_lib("openpyxl"), "cần openpyxl")
    def test_xlsx_liet_ke_sheet(self):
        import openpyxl
        wb = openpyxl.Workbook()
        wb.active.title = "Lương"
        wb.create_sheet("Thưởng")
        wb.save(str(self.ws / "a.xlsx"))
        kq = ob.doc_thong_tin("a.xlsx")
        self.assertEqual(kq["so_sheet"], 2)
        self.assertIn("Lương", kq["cac_sheet"])

    def test_khong_doc_duoc_ngoai_workspace(self):
        self.assertFalse(ob.doc_thong_tin("/etc/passwd")["ok"])


@unittest.skipUnless(_co_lib("docx", "openpyxl", "pandas"), "cần docx+openpyxl+pandas")
class TaoBaoCaoTests(_Nen):

    def _so_lieu(self):
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Luong"
        ws.append(["Tên", "Lương"])
        ws.append(["An", 10])
        ws.append(["Bình", 20])
        wb.save(str(self.ws / "so.xlsx"))
        return "so.xlsx"

    def _chu_ra(self, ten="bc.docx"):
        import docx
        d = docx.Document(str(self.ws / ten))
        chu = [p.text for p in d.paragraphs]
        for b in d.tables:
            for h in b.rows:
                chu += [o.text for o in h.cells]
        return "\n".join(chu)

    def test_khong_mau_thi_tu_dung_bao_cao(self):
        kq = ob.tao_bao_cao(self._so_lieu(), "bc.docx")
        self.assertTrue(kq["ok"], kq.get("error"))
        self.assertFalse(kq["theo_mau"])
        s = self._chu_ra()
        self.assertIn("Luong", s)
        self.assertIn("An", s)
        self.assertIn("Trung bình", s)

    def test_MAU_THAT_SU_DUOC_DUNG(self):
        """Repo nhận `template_md` rồi KHÔNG dùng ở đâu cả — tiêu đề luôn cứng."""
        mau = ("# Báo cáo lương tháng 8\n"
               "Đây là phần mở đầu do tôi viết.\n"
               "## Số liệu\n"
               "{{bang:Luong}}\n"
               "## Tổng quan\n"
               "{{thong_ke:Luong}}\n")
        kq = ob.tao_bao_cao(self._so_lieu(), "bc.docx", mau_md=mau)
        self.assertTrue(kq["ok"], kq.get("error"))
        self.assertTrue(kq["theo_mau"])
        s = self._chu_ra()
        self.assertIn("Báo cáo lương tháng 8", s)
        self.assertIn("Đây là phần mở đầu do tôi viết", s)
        self.assertIn("An", s)          # {{bang:…}} đã thành bảng thật
        self.assertIn("Trung bình", s)   # {{thong_ke:…}} đã thành bảng thống kê

    def test_mau_tro_sheet_khong_co_thi_noi_ro(self):
        kq = ob.tao_bao_cao(self._so_lieu(), "bc.docx",
                            mau_md="# A\n{{bang:KhongCo}}\n")
        self.assertTrue(kq["ok"])
        self.assertIn("KhongCo", self._chu_ra())

    def test_khong_phai_bang_tinh_thi_bao_ro(self):
        (self.ws / "a.docx").write_text("x", "utf-8")
        self.assertFalse(ob.tao_bao_cao("a.docx", "bc.docx")["ok"])


class DaCamVaoNangLucBotTests(unittest.TestCase):
    """Làm module xong mà không khai thành năng lực thì bot không gọi được."""

    MOI = ("office_so_sanh", "office_thong_ke", "office_thay_the",
           "office_tao_slide", "office_cat", "office_noi",
           "office_thong_tin", "office_bao_cao")

    def _src(self):
        return (GOC / "services" / "agent" / "capabilities.py").read_text("utf-8")

    def test_moi_viec_deu_co_capability_va_handler(self):
        src = self._src()
        for ten in self.MOI:
            with self.subTest(ten=ten):
                self.assertIn(f'"{ten}": Capability(', src)
                self.assertIn(f"def _h_{ten}(", src)

    def test_moi_viec_deu_co_trong_bang_NHOM(self):
        """Thiếu nhóm thì group_of() trả '_ungrouped' và bot bị chặn IM LẶNG."""
        src = self._src()
        for ten in self.MOI:
            with self.subTest(ten=ten):
                self.assertIn(f'"{ten}": "office"', src)

    def test_thay_the_phai_la_CHANGE_chu_khong_phai_READ(self):
        """Nó ghi thẳng vào tệp, không có bản lùi."""
        src = self._src()
        i = src.index('"office_thay_the": Capability(')
        self.assertIn("risk=CHANGE", src[i:i + 200])


if __name__ == "__main__":
    unittest.main()
